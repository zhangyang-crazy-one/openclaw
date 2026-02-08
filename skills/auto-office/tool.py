#!/usr/bin/env python3
"""
Auto-Office Tool - Excel & PPT 生成器
"""

import argparse
import json
import csv
import sys
from pathlib import Path
from datetime import datetime

# Excel 生成
try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    from openpyxl.utils import get_column_letter
except ImportError:
    print("❌ 请安装: pip install openpyxl xlsxwriter pandas")
    sys.exit(1)

# PPT 生成
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError as e:
    print(f"❌ 请安装: pip install python-pptx ({e})")
    sys.exit(1)


class ExcelGenerator:
    """Excel 文件生成器"""

    def __init__(self, output_path):
        self.output_path = Path(output_path)
        self.wb = Workbook()
        self.ws = self.wb.active
        self.current_row = 1

        # 样式
        self.header_font = Font(bold=True, size=12)
        self.header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        self.header_font_white = Font(bold=True, size=12, color="FFFFFF")
        self.alignment = Alignment(horizontal="center", vertical="center")
        self.border = Border(
            left=Side(style='thin'),
            right=Side(style='thin'),
            top=Side(style='thin'),
            bottom=Side(style='thin')
        )

    def add_headers(self, headers, bold=True, fill_color=None):
        """添加表头"""
        for col, header in enumerate(headers, 1):
            cell = self.ws.cell(row=self.current_row, column=col, value=header)
            cell.font = self.header_font_white if fill_color else self.header_font
            if fill_color:
                cell.fill = PatternFill(start_color=fill_color, end_color=fill_color, fill_type="solid")
            cell.alignment = self.alignment
            cell.border = self.border
        self.current_row += 1

    def add_row(self, row_data):
        """添加数据行"""
        for col, value in enumerate(row_data, 1):
            cell = self.ws.cell(row=self.current_row, column=col, value=value)
            cell.alignment = self.alignment
            cell.border = self.border
        self.current_row += 1

    def add_data_from_csv(self, csv_data):
        """从CSV数据添加多行"""
        lines = csv_data.strip().split('\n')
        reader = csv.reader(lines)
        for row in reader:
            self.add_row(row)

    def add_data_from_json(self, json_data):
        """从JSON数据添加"""
        data = json.loads(json_data)
        for row in data:
            self.add_row(row)

    def set_bold_header(self, fill_color="4472C4"):
        """设置表头样式"""
        # 样式已在 add_headers 中应用
        pass

    def auto_width(self):
        """自动列宽"""
        for column in self.ws.columns:
            max_length = 0
            column_letter = get_column_letter(column[0].column)
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = min(max_length + 2, 50)
            self.ws.column_dimensions[column_letter].width = adjusted_width

    def freeze_header(self):
        """冻结首行"""
        self.ws.freeze_panes = 'A2'

    def save(self):
        """保存文件"""
        self.wb.save(str(self.output_path))
        print(f"✅ Excel 已保存: {self.output_path}")
        return str(self.output_path)


class PPTGenerator:
    """PPT 文件生成器"""

    def __init__(self, output_path, title="演示文稿"):
        self.output_path = Path(output_path)
        self.prs = Presentation()
        self.title = title

        # 主题颜色
        self.themes = {
            'simple': {'primary': '2F5496', 'accent': '4472C4'},
            'modern': {'primary': '1F4E79', 'accent': '5B9BD5'},
            'professional': {'primary': '363636', 'accent': '538135'}
        }

    def set_theme(self, theme_name='simple'):
        """设置主题"""
        self.current_theme = self.themes.get(theme_name, self.themes['simple'])

    def add_title_slide(self, title, subtitle=""):
        """添加标题页"""
        slide_layout = self.prs.slide_layouts[0]  # Title Slide
        slide = self.prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title
        if subtitle:
            subtitle_shape.text = subtitle

    def add_content_slide(self, title, content_lines):
        """添加内容页"""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        for i, line in enumerate(content_lines):
            if i == 0:
                tf.text = line
            else:
                p = tf.add_paragraph()
                p.text = line

    def add_slide(self, title, content="", layout_idx=1):
        """添加通用幻灯片"""
        slide_layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if content:
            body = slide.placeholders[1]
            body.text = content
        return slide

    def add_chart(self, data, labels, title="图表", slide_idx=None):
        """添加图表"""
        # 创建临时Excel生成图表
        from openpyxl import Workbook
        from openpyxl.chart import BarChart, Reference

        wb = Workbook()
        ws = wb.active
        ws.title = "ChartData"

        # 添加数据
        ws['A1'] = ''
        for i, label in enumerate(labels, 1):
            ws.cell(row=1, column=i+1, value=label)
        ws['A2'] = '数值'
        for i, value in enumerate(data, 1):
            ws.cell(row=2, column=i+1, value=value)

        # 添加到PPT
        if slide_idx is not None:
            slide = self.prs.slides[slide_idx]
        else:
            slide_layout = self.prs.slide_layouts[5]  # Blank
            slide = self.prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title

        # 简化：添加文本框显示数据
        left = Inches(2)
        top = Inches(2)
        width = Inches(6)
        height = Inches(4)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"{title}\n\n"
        for label, value in zip(labels, data):
            tf.text += f"• {label}: {value}\n"

    def save(self):
        """保存文件"""
        self.prs.save(str(self.output_path))
        print(f"✅ PPT 已保存: {self.output_path}")
        return str(self.output_path)


def main():
    parser = argparse.ArgumentParser(description="Auto-Office: Excel & PPT 生成器")
    subparsers = parser.add_subparsers(dest='command', help='子命令')

    # Excel 命令
    excel_parser = subparsers.add_parser('excel', help='创建Excel文件')
    excel_parser.add_argument('-o', '--output', required=True, help='输出文件路径')
    excel_parser.add_argument('-d', '--data', help='数据（JSON或CSV格式）')
    excel_parser.add_argument('--headers', help='表头（逗号分隔）')
    excel_parser.add_argument('--sheet', default='Sheet1', help='工作表名')
    excel_parser.add_argument('--bold-header', action='store_true', help='表头加粗')
    excel_parser.add_argument('--auto-width', action='store_true', help='自动列宽')
    excel_parser.add_argument('--freeze-header', action='store_true', help='冻结首行')

    # PPT 命令
    ppt_parser = subparsers.add_parser('ppt', help='创建PPT文件')
    ppt_parser.add_argument('-o', '--output', required=True, help='输出文件路径')
    ppt_parser.add_argument('-t', '--title', default='演示文稿', help='演示文稿标题')
    ppt_parser.add_argument('-S', '--slides', help='幻灯片列表（格式: "标题:内容"）')
    ppt_parser.add_argument('--theme', default='simple', choices=['simple', 'modern', 'professional'], help='主题')
    ppt_parser.add_argument('--add-chart', action='store_true', help='添加示例图表')

    args = parser.parse_args()

    if args.command == 'excel':
        excel = ExcelGenerator(args.output)

        # 添加表头
        if args.headers:
            headers = [h.strip() for h in args.headers.split(',')]
            excel.add_headers(headers, bold=args.bold_header)

        # 添加数据
        if args.data:
            if args.data.startswith('['):
                excel.add_data_from_json(args.data)
            else:
                excel.add_data_from_csv(args.data)

        # 其他选项
        if args.auto_width:
            excel.auto_width()
        if args.freeze_header:
            excel.freeze_header()

        result = excel.save()
        print(f"📊 Excel: {result}")

    elif args.command == 'ppt':
        ppt = PPTGenerator(args.output, args.title)
        ppt.set_theme(args.theme)

        # 解析幻灯片
        if args.slides:
            for slide_info in args.slides.split('|'):
                if ':' in slide_info:
                    title, content = slide_info.split(':', 1)
                    ppt.add_slide(title.strip(), content.strip())
                else:
                    ppt.add_slide(slide_info.strip())

        if args.add_chart:
            ppt.add_chart([10, 20, 30, 25], ['Q1', 'Q2', 'Q3', 'Q4'], "季度数据")

        result = ppt.save()
        print(f"📑 PPT: {result}")

    else:
        parser.print_help()


if __name__ == '__main__':
    main()
